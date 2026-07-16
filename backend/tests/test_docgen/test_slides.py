"""Slide content infrastructure tests: typography, geometry, flattening,
autofit, and consulting-grade PPTX render round trips."""

from __future__ import annotations

from typing import TYPE_CHECKING

import pytest

if TYPE_CHECKING:
    from pathlib import Path

from leagent.docgen.model import DeckSpec
from leagent.docgen.slides import (
    EMU_PER_IN,
    BulletPara,
    DeckTypography,
    Region,
    SlideGeometry,
    fit_body_size,
    flatten_body,
    is_dark_color,
)
from leagent.docgen.themes import get_theme

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# 1x1 transparent PNG.
_PX = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8"
    "z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _geom() -> SlideGeometry:
    return SlideGeometry(slide_w=int(13.333 * EMU_PER_IN), slide_h=int(7.5 * EMU_PER_IN))


# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------


def test_typography_scale_from_theme() -> None:
    theme = get_theme("midnight_executive", kind="deck")
    typo = DeckTypography.from_theme(theme)

    assert typo.display.size == theme.deck.title_size
    assert typo.slide_title.size == theme.deck.slide_title_size
    assert typo.kicker.uppercase and typo.kicker.letter_spacing_pt > 0
    # Bullet levels shrink monotonically and never go below 10pt.
    sizes = [typo.level(i).text.size for i in range(5)]
    assert sizes == sorted(sizes, reverse=True)
    assert all(s >= 10.0 for s in sizes)
    # Deeper levels indent further.
    assert typo.level(2).indent_in > typo.level(0).indent_in
    # Out-of-range levels clamp instead of raising.
    assert typo.level(99) == typo.level(4)


def test_is_dark_color() -> None:
    assert is_dark_color("#0B1B33") is True
    assert is_dark_color("#FFFFFF") is False
    assert is_dark_color("#F4F6FA") is False
    assert is_dark_color("garbage") is True  # unparseable -> treated as black


# ---------------------------------------------------------------------------
# Geometry
# ---------------------------------------------------------------------------


def test_content_top_accounts_for_header_parts() -> None:
    geom = _geom()
    bare = geom.content_top(has_title=False, has_kicker=False, has_subtitle=False)
    titled = geom.content_top(has_title=True, has_kicker=False, has_subtitle=False)
    kicked = geom.content_top(has_title=True, has_kicker=True, has_subtitle=False)
    full = geom.content_top(has_title=True, has_kicker=True, has_subtitle=True)
    assert bare == geom.margin
    assert bare < titled < kicked < full


def test_takeaway_carves_content_region() -> None:
    geom = _geom()
    top = geom.content_top(has_title=True, has_kicker=False, has_subtitle=False)
    plain = geom.content_region(top)
    carved = geom.content_region(top, has_takeaway=True)
    assert carved.height < plain.height
    takeaway = geom.takeaway_region()
    assert carved.bottom <= takeaway.top
    assert takeaway.bottom <= geom.slide_h - geom.footer_h + 1


@pytest.mark.parametrize("side", ["left", "right", "top"])
def test_split_partitions_with_gutter(side: str) -> None:
    geom = _geom()
    region = geom.content_region(geom.margin)
    text, media = geom.split(region, 0.4, side=side)
    if side == "top":
        assert media.top == region.top and text.top > media.bottom
        assert text.height + media.height + geom.gutter == region.height
    else:
        assert text.width + media.width + geom.gutter == region.width
        if side == "left":
            assert media.left == region.left and text.left > media.right
        else:
            assert text.left == region.left and media.left > text.right
    # Ratio is clamped to sane bounds.
    _, huge = geom.split(region, 5.0)
    assert huge.width <= int((region.width - geom.gutter) * 0.8)


def test_columns_equal_and_weighted() -> None:
    geom = _geom()
    region = geom.content_region(geom.margin)
    cols = geom.columns(region, 3)
    assert len(cols) == 3
    assert cols[0].width == cols[1].width == cols[2].width
    assert cols[1].left == cols[0].right + geom.gutter

    weighted = geom.columns(region, 2, weights=[2.0, 1.0])
    assert weighted[0].width > weighted[1].width


# ---------------------------------------------------------------------------
# Body flattening
# ---------------------------------------------------------------------------


def test_flatten_body_multi_level_content() -> None:
    md = (
        "#### 重点\n\n"
        "- 一级要点\n"
        "  - 二级要点\n"
        "    - 三级要点\n"
        "1. 第一步\n"
        "2. 第二步\n"
        "- [x] 已完成\n"
        "- [ ] 待办\n\n"
        "普通段落。\n\n"
        "> 引用语\n\n"
        "```\ncode line\n```\n"
    )
    paras = flatten_body(md)
    kinds = [(p.kind, p.level) for p in paras]
    assert ("heading", 0) in kinds
    assert ("bullet", 0) in kinds and ("bullet", 1) in kinds and ("bullet", 2) in kinds
    numbered = [p for p in paras if p.kind == "numbered"]
    assert [p.number for p in numbered] == [1, 2]
    checked = [p.checked for p in paras if p.checked is not None]
    assert checked == [True, False]
    assert any(p.kind == "para" for p in paras)
    assert any(p.kind == "quote" for p in paras)
    assert any(p.kind == "code" for p in paras)


def test_flatten_body_empty() -> None:
    assert flatten_body("") == []
    assert flatten_body("   \n  ") == []


# ---------------------------------------------------------------------------
# Autofit
# ---------------------------------------------------------------------------


def test_fit_body_size_shrinks_dense_content() -> None:
    theme = get_theme("professional", kind="deck")
    typo = DeckTypography.from_theme(theme)
    region = Region(0, 0, int(6 * EMU_PER_IN), int(4 * EMU_PER_IN))

    short = [BulletPara(text="one point", kind="bullet")]
    assert fit_body_size(short, typo, region) == 1.0

    dense = [
        BulletPara(text="这是一条相当长的中文要点内容,用来测试自动缩放" * 3, kind="bullet")
        for _ in range(14)
    ]
    scale = fit_body_size(dense, typo, region)
    assert scale < 1.0
    assert scale >= 0.62  # floor


# ---------------------------------------------------------------------------
# Render round trips (python-pptx)
# ---------------------------------------------------------------------------


def _all_text(prs) -> str:
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def test_render_pptx_consulting_anatomy(tmp_path: Path) -> None:
    """kicker / takeaway / columns / numbered bullets survive the render."""
    pptx = pytest.importorskip("pptx")
    from leagent.docgen.renderers.pptx import render_pptx

    spec = DeckSpec.model_validate(
        {
            "title": "增长战略",
            "slides": [
                {"layout": "title", "title": "增长战略"},
                {
                    "layout": "content",
                    "kicker": "market analysis",
                    "title": "华东市场三年翻倍",
                    "body": "- 一级要点\n  - 二级要点\n1. 第一步\n2. 第二步",
                    "takeaway": "渠道布局须在 Q3 前完成",
                },
                {
                    "layout": "columns",
                    "title": "三条路径",
                    "columns": [
                        {"heading": "自建", "body": "- 投入高"},
                        {"heading": "联营", "body": "- 平衡", "emphasis": True},
                        {"heading": "代理", "body": "- 风险高"},
                    ],
                },
            ],
        }
    )
    out = tmp_path / "deck.pptx"
    result = render_pptx(spec, out)
    assert result["success"] is True
    assert result["warnings"] == []

    prs = pptx.Presentation(str(out))
    joined = _all_text(prs)
    # Kicker renders uppercase.
    assert "MARKET ANALYSIS" in joined
    assert "渠道布局须在 Q3 前完成" in joined
    for heading in ("自建", "联营", "代理"):
        assert heading in joined

    # Numbered items use native auto-numbering; leveled bullets carry
    # hanging-indent geometry.
    content_xml = prs.slides[1].shapes._spTree.xml
    assert "buAutoNum" in content_xml
    assert 'indent="-' in content_xml

    # The emphasized column draws a card (an autoshape besides text boxes).
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    col_slide_shapes = [s.shape_type for s in prs.slides[2].shapes]
    assert MSO_SHAPE_TYPE.AUTO_SHAPE in col_slide_shapes


def test_render_pptx_backgrounds(tmp_path: Path) -> None:
    """Solid / gradient / image backgrounds apply, and palettes adapt."""
    pptx = pytest.importorskip("pptx")
    from leagent.docgen.renderers.pptx import render_pptx

    spec = DeckSpec.model_validate(
        {
            "title": "背景测试",
            "theme": "professional",  # light theme
            "slides": [
                {"layout": "title", "title": "背景测试"},
                {
                    "layout": "content",
                    "title": "深色纯色",
                    "body": "- 要点",
                    "background": {"color": "#0B1B33"},
                },
                {
                    "layout": "content",
                    "title": "渐变",
                    "body": "- 要点",
                    "background": {"gradient": ["#0B1B33", "#1E3A5F"]},
                },
                {
                    "layout": "content",
                    "title": "图片背景",
                    "body": "- 要点",
                    "image": {"base64_data": _PX, "position": "background"},
                },
            ],
        }
    )
    out = tmp_path / "deck.pptx"
    result = render_pptx(spec, out)
    assert result["success"] is True

    prs = pptx.Presentation(str(out))
    # Slide 2: solid dark background element present.
    bg1 = prs.slides[1].background.fill
    assert bg1.fore_color.rgb == pptx.dml.color.RGBColor.from_string("0B1B33")
    # On the dark slide, body text is not the light theme's dark default.
    dark_runs = [
        run
        for shape in prs.slides[1].shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
        if "要点" in run.text
    ]
    assert dark_runs
    assert str(dark_runs[0].font.color.rgb) == "F0F2F5"
    # Slide 3: gradient fill XML present.
    assert "gradFill" in prs.slides[2].background._cSld.xml
    # Slide 4: background image placed as a picture shape.
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    assert any(
        s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in prs.slides[3].shapes
    )


def test_render_pptx_image_text_split(tmp_path: Path) -> None:
    """image.position right/left/top places media beside autofit body text."""
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from leagent.docgen.renderers.pptx import render_pptx

    spec = DeckSpec.model_validate(
        {
            "title": "图文",
            "slides": [
                {"layout": "title", "title": "图文"},
                {
                    "layout": "content",
                    "title": "右图",
                    "body": "- 文字在左",
                    "image": {"base64_data": _PX, "position": "right", "ratio": 0.4},
                },
                {
                    "layout": "image",
                    "title": "左图",
                    "body": "- 文字在右",
                    "image": {"base64_data": _PX, "position": "left", "ratio": 0.35},
                },
            ],
        }
    )
    out = tmp_path / "deck.pptx"
    result = render_pptx(spec, out)
    assert result["success"] is True
    assert result["content_stats"]["images"] == 2
    assert result["warnings"] == []

    prs = pptx.Presentation(str(out))
    for slide_idx, text_side in ((1, "left"), (2, "right")):
        slide = prs.slides[slide_idx]
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1
        body_boxes = [
            s
            for s in slide.shapes
            if s.has_text_frame and "文字" in s.text_frame.text
        ]
        assert len(body_boxes) == 1
        if text_side == "left":
            assert body_boxes[0].left < pics[0].left
        else:
            assert body_boxes[0].left > pics[0].left


def test_render_pptx_autofit_never_overflows_region(tmp_path: Path) -> None:
    """Dense body text shrinks instead of spilling past the footer."""
    pptx = pytest.importorskip("pptx")
    from leagent.docgen.renderers.pptx import render_pptx

    dense_body = "\n".join(
        f"- 第{i}条要点:这是一条比较长的说明文字,覆盖自动缩放逻辑的验证场景" for i in range(1, 15)
    )
    spec = DeckSpec.model_validate(
        {
            "title": "密度",
            "slides": [
                {"layout": "title", "title": "密度"},
                {"layout": "content", "title": "高密度", "body": dense_body},
                {"layout": "content", "title": "低密度", "body": "- 只有一条"},
            ],
        }
    )
    out = tmp_path / "deck.pptx"
    assert render_pptx(spec, out)["success"] is True

    prs = pptx.Presentation(str(out))

    def _body_sizes(slide) -> list[int]:
        return [
            run.font.size.pt
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
            if "要点" in run.text or "只有一条" in run.text
        ]

    dense_sizes = _body_sizes(prs.slides[1])
    sparse_sizes = _body_sizes(prs.slides[2])
    assert dense_sizes and sparse_sizes
    assert max(dense_sizes) < max(sparse_sizes)


def test_render_pptx_background_image_with_overlay_only(tmp_path: Path) -> None:
    """overlay-only background must not block image.position=background."""
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from leagent.docgen.renderers.pptx import render_pptx

    spec = DeckSpec.model_validate(
        {
            "title": "封面",
            "slides": [
                {
                    "layout": "title",
                    "title": "LeAgent",
                    "subtitle": "AI 助手",
                    "image": {"base64_data": _PX, "position": "background"},
                    "background": {"overlay": 0.35},
                },
                {
                    "layout": "closing",
                    "title": "Thank You",
                    "image": {"base64_data": _PX, "position": "background"},
                    "background": {"overlay": 0.25},
                },
            ],
        }
    )
    out = tmp_path / "overlay.pptx"
    result = render_pptx(spec, out)
    assert result["success"] is True
    assert result["content_stats"]["images"] == 2
    assert result["warnings"] == []

    prs = pptx.Presentation(str(out))
    for slide in prs.slides:
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(pics) == 1


def test_render_pptx_two_column_with_side_image(tmp_path: Path) -> None:
    """two_column + image.position=right places media beside left text."""
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from leagent.docgen.renderers.pptx import render_pptx

    spec = DeckSpec.model_validate(
        {
            "title": "关于",
            "slides": [
                {"layout": "title", "title": "关于"},
                {
                    "layout": "two_column",
                    "title": "产品亮点",
                    "left": "- 智能推理\n- 专业工具",
                    "image": {"base64_data": _PX, "position": "right", "ratio": 0.45},
                    "takeaway": "不止于聊天",
                },
            ],
        }
    )
    out = tmp_path / "two_col.pptx"
    result = render_pptx(spec, out)
    assert result["success"] is True
    assert result["content_stats"]["images"] == 1
    assert result["warnings"] == []

    prs = pptx.Presentation(str(out))
    slide = prs.slides[1]
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1
    body_boxes = [
        s for s in slide.shapes if s.has_text_frame and "智能推理" in s.text_frame.text
    ]
    assert len(body_boxes) == 1
    assert body_boxes[0].left < pics[0].left
