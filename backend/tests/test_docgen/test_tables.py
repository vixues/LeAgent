"""Shared table engine: normalization, inference, semantics, style contract."""

from __future__ import annotations

from leagent.docgen.model import TableBlock
from leagent.docgen.tables import process_table, resolve_table_style
from leagent.docgen.themes import get_theme


def _block(**kwargs) -> TableBlock:
    return TableBlock.model_validate({"type": "table", **kwargs})


# ---------------------------------------------------------------------------
# Normalization
# ---------------------------------------------------------------------------


def test_ragged_rows_are_padded_to_uniform_grid() -> None:
    pt = process_table(_block(columns=["A", "B", "C"], rows=[["1"], ["1", "2", "3", "4"]]))
    assert pt.col_count == 4  # widest row wins
    assert all(len(r) == 4 for r in pt.body)
    assert len(pt.header) == 4


def test_first_row_is_header_when_columns_omitted() -> None:
    pt = process_table(_block(rows=[["名称", "数量"], ["苹果", "3"]]))
    assert [c.text for c in pt.header] == ["名称", "数量"]
    assert len(pt.body) == 1


def test_empty_table() -> None:
    pt = process_table(_block(rows=[]))
    assert not pt.header
    assert not pt.body


# ---------------------------------------------------------------------------
# Column inference + alignment
# ---------------------------------------------------------------------------


def test_numeric_columns_right_align_text_left() -> None:
    pt = process_table(
        _block(
            columns=["项目", "金额", "占比"],
            rows=[["研发", "1,200", "34%"], ["销售", "800", "22%"], ["行政", "410", "12%"]],
        )
    )
    assert [c.align for c in pt.columns] == ["left", "right", "right"]
    assert pt.columns[1].kind == "number"
    assert pt.columns[2].kind == "percent"


def test_explicit_align_overrides_inference() -> None:
    pt = process_table(
        _block(
            columns=["项目", "金额"],
            rows=[["研发", "1200"], ["销售", "800"]],
            align=["center", "left"],
        )
    )
    assert [c.align for c in pt.columns] == ["center", "left"]


def test_currency_and_date_kinds() -> None:
    pt = process_table(
        _block(
            columns=["期间", "收入"],
            rows=[["2025-01", "¥1,200"], ["2025-02", "¥1,350"], ["2025-03", "$900"]],
        )
    )
    assert pt.columns[0].kind == "date"
    assert pt.columns[0].align == "center"
    assert pt.columns[1].kind == "currency"
    assert pt.columns[1].align == "right"


# ---------------------------------------------------------------------------
# Number polish
# ---------------------------------------------------------------------------


def test_bare_integers_gain_thousands_separators() -> None:
    pt = process_table(_block(columns=["k", "v"], rows=[["a", "1234567"], ["b", "89012.5"]]))
    assert pt.body[0][1].text == "1,234,567"
    assert pt.body[1][1].text == "89,012.5"


def test_number_format_can_be_disabled_and_skips_non_bare_cells() -> None:
    pt = process_table(
        _block(columns=["k", "v"], rows=[["a", "1234567"]], number_format=False)
    )
    assert pt.body[0][1].text == "1234567"
    pt2 = process_table(_block(columns=["k", "v"], rows=[["a", "约12345元"], ["b", "1,234"]]))
    assert pt2.body[0][1].text == "约12345元"  # untouched
    assert pt2.body[1][1].text == "1,234"  # already separated


# ---------------------------------------------------------------------------
# Total row + delta semantics
# ---------------------------------------------------------------------------


def test_total_row_auto_detected_and_bold() -> None:
    pt = process_table(
        _block(
            columns=["项目", "金额"],
            rows=[["研发", "1200"], ["销售", "800"], ["合计", "2000"]],
        )
    )
    assert pt.total_row_index == 2
    assert all(cell.bold for cell in pt.body[2])
    # Total label must not poison column-kind inference.
    assert pt.columns[0].kind == "text"


def test_total_row_flag_overrides_detection() -> None:
    rows = [["r1", "1"], ["r2", "2"]]
    forced = process_table(_block(columns=["k", "v"], rows=rows, total_row=True))
    assert forced.total_row_index == 1
    suppressed = process_table(
        _block(columns=["k", "v"], rows=[*rows, ["Total", "3"]], total_row=False)
    )
    assert suppressed.total_row_index is None


def test_delta_cells_carry_polarity() -> None:
    pt = process_table(
        _block(
            columns=["指标", "同比"],
            rows=[["收入", "+8.2%"], ["成本", "-3.1%"], ["利润", "(1,200)"]],
        )
    )
    assert pt.body[0][1].polarity == "positive"
    assert pt.body[1][1].polarity == "negative"
    assert pt.body[2][1].polarity == "negative"
    assert pt.body[0][0].polarity is None


# ---------------------------------------------------------------------------
# Widths
# ---------------------------------------------------------------------------


def test_width_fractions_sum_to_one_and_favor_long_columns() -> None:
    pt = process_table(
        _block(
            columns=["ID", "描述"],
            rows=[["1", "这是一个非常长的中文描述，用于测试列宽分配的效果"], ["2", "短"]],
        )
    )
    fractions = pt.width_fractions()
    assert abs(sum(fractions) - 1.0) < 1e-6
    assert fractions[1] > fractions[0]


def test_explicit_widths_win() -> None:
    pt = process_table(
        _block(columns=["a", "b"], rows=[["1", "2"]], widths=[30, 70])
    )
    assert pt.width_fractions() == [0.3, 0.7]


# ---------------------------------------------------------------------------
# Style contract
# ---------------------------------------------------------------------------


def test_style_variants_resolve_distinct_contracts() -> None:
    theme = get_theme("professional", kind="document")
    default = resolve_table_style(theme, "default")
    minimal = resolve_table_style(theme, "minimal")
    grid = resolve_table_style(theme, "grid")

    assert default.header_fill == theme.colors.primary
    assert default.header_text == "#FFFFFF"
    assert minimal.header_fill is None  # open header, rules only
    assert minimal.outer_rule is not None
    assert grid.grid is True


def test_dark_deck_style_keeps_contrast() -> None:
    theme = get_theme("midnight_executive", kind="deck")
    ts = resolve_table_style(theme, "default", dark=True)
    assert ts.header_fill == theme.colors.accent
    assert ts.header_text == theme.colors.background


def test_zebra_suppressed_for_minimal_and_overridable() -> None:
    theme = get_theme("professional", kind="document")
    minimal = process_table(
        _block(columns=["a"], rows=[["1"], ["2"]], style="minimal"), theme=theme
    )
    assert minimal.zebra is False
    forced_off = process_table(
        _block(columns=["a"], rows=[["1"], ["2"]], zebra=False), theme=theme
    )
    assert forced_off.zebra is False


# ---------------------------------------------------------------------------
# Render round trips (the engine drives every format identically)
# ---------------------------------------------------------------------------

_FINANCIAL_TABLE = {
    "type": "table",
    "columns": ["项目", "金额", "同比"],
    "rows": [
        ["研发", "1200000", "+8.2%"],
        ["销售", "800000", "-3.1%"],
        ["合计", "2000000", "+3.4%"],
    ],
    "caption": "费用一览",
}


def _doc_spec():
    from leagent.docgen.model import DocumentSpec

    return DocumentSpec.model_validate(
        {"title": "表格测试", "blocks": [_FINANCIAL_TABLE]}
    )


def test_html_render_applies_table_semantics(tmp_path) -> None:
    from leagent.docgen.renderers.html import render_html

    out = tmp_path / "t.html"
    assert render_html(_doc_spec(), out)["success"]
    doc = out.read_text(encoding="utf-8")
    assert "1,200,000" in doc                # number polish
    assert "text-align:right" in doc         # numeric auto-alignment
    assert "font-weight:700" in doc          # total row emphasis
    assert "#1E8449" in doc and "#C0392B" in doc  # delta polarity colors
    assert "费用一览" in doc


def test_markdown_render_emits_alignment_markers(tmp_path) -> None:
    from leagent.docgen.renderers.html import render_markdown

    out = tmp_path / "t.md"
    assert render_markdown(_doc_spec(), out)["success"]
    doc = out.read_text(encoding="utf-8")
    assert "| ---  | ---: | ---: |".replace("  ", " ") in doc.replace("  ", " ")
    assert "| 合计 | 2,000,000 | +3.4% |" in doc


def test_docx_render_applies_table_semantics(tmp_path) -> None:
    import pytest

    docx = pytest.importorskip("docx")
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    from leagent.docgen.renderers.docx import render_docx

    out = tmp_path / "t.docx"
    assert render_docx(_doc_spec(), out)["success"]
    document = docx.Document(str(out))
    table = document.tables[0]
    # Numeric column right-aligned.
    amount_cell = table.cell(1, 1)
    assert amount_cell.paragraphs[0].alignment == WD_ALIGN_PARAGRAPH.RIGHT
    assert "1,200,000" in amount_cell.text
    # Total row bold.
    total_runs = [r for p in table.cell(3, 0).paragraphs for r in p.runs]
    assert total_runs and all(r.font.bold for r in total_runs)


def test_pptx_render_table_slide(tmp_path) -> None:
    import pytest

    pptx = pytest.importorskip("pptx")
    from pptx.enum.text import PP_ALIGN

    from leagent.docgen.model import DeckSpec
    from leagent.docgen.renderers.pptx import render_pptx

    deck = DeckSpec.model_validate(
        {
            "title": "表格",
            "slides": [{"layout": "table", "title": "费用", "table": _FINANCIAL_TABLE}],
        }
    )
    out = tmp_path / "t.pptx"
    assert render_pptx(deck, out)["success"]
    prs = pptx.Presentation(str(out))
    tables = [sh.table for sl in prs.slides for sh in sl.shapes if sh.has_table]
    assert tables
    table = tables[0]
    assert "1,200,000" in table.cell(1, 1).text
    assert table.cell(1, 1).text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT


def test_pdf_render_with_financial_table(tmp_path) -> None:
    import pytest

    pytest.importorskip("reportlab")
    from leagent.docgen.renderers.pdf import render_pdf

    out = tmp_path / "t.pdf"
    result = render_pdf(_doc_spec(), out)
    assert result["success"] is True
    assert result["content_stats"]["tables"] == 1
    assert out.stat().st_size > 0
