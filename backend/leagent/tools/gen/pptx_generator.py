"""PowerPoint Generator Tool - Create .pptx presentations programmatically.

Uses python-pptx for creating presentations with slides, layouts, images, tables, and charts.
"""

from __future__ import annotations

import base64
import io
from pathlib import Path
from typing import Any

import structlog

from leagent.tools.base import SyncTool, ToolCategory, ToolContext

logger = structlog.get_logger(__name__)

THEME_PALETTES: dict[str, dict[str, str]] = {
    "midnight_executive": {
        "primary": "1E2761",
        "secondary": "CADCFC",
        "accent": "FFFFFF",
        "header_font": "Georgia",
        "body_font": "Calibri",
    },
    "forest_moss": {
        "primary": "2C5F2D",
        "secondary": "97BC62",
        "accent": "F5F5F5",
        "header_font": "Cambria",
        "body_font": "Calibri",
    },
    "coral_energy": {
        "primary": "F96167",
        "secondary": "F9E795",
        "accent": "2F3C7E",
        "header_font": "Arial Black",
        "body_font": "Arial",
    },
    "warm_terracotta": {
        "primary": "B85042",
        "secondary": "E7E8D1",
        "accent": "A7BEAE",
        "header_font": "Palatino",
        "body_font": "Calibri",
    },
    "ocean_gradient": {
        "primary": "065A82",
        "secondary": "1C7293",
        "accent": "21295C",
        "header_font": "Trebuchet MS",
        "body_font": "Calibri",
    },
    "charcoal_minimal": {
        "primary": "36454F",
        "secondary": "F2F2F2",
        "accent": "212121",
        "header_font": "Calibri",
        "body_font": "Calibri Light",
    },
    "teal_trust": {
        "primary": "028090",
        "secondary": "00A896",
        "accent": "02C39A",
        "header_font": "Georgia",
        "body_font": "Calibri",
    },
    "berry_cream": {
        "primary": "6D2E46",
        "secondary": "A26769",
        "accent": "ECE2D0",
        "header_font": "Palatino",
        "body_font": "Calibri",
    },
    "sage_calm": {
        "primary": "84B59F",
        "secondary": "69A297",
        "accent": "50808E",
        "header_font": "Georgia",
        "body_font": "Calibri",
    },
    "cherry_bold": {
        "primary": "990011",
        "secondary": "FCF6F5",
        "accent": "2F3C7E",
        "header_font": "Impact",
        "body_font": "Arial",
    },
}

TYPOGRAPHY_SIZES: dict[str, int] = {
    "slide_title": 36,
    "section_header": 24,
    "body": 16,
    "caption": 11,
}


class PptxGeneratorTool(SyncTool):
    """Generate PowerPoint presentations (.pptx) with rich slide content.

    Features:
    - Create presentations from scratch or templates
    - Multiple slide layouts (title, content, two-column, blank, etc.)
    - 10 curated theme palettes with matching typography
    - Add text, images, tables, charts, and shapes
    - Gradient and solid slide backgrounds
    - Automatic slide numbering
    - Speaker notes per slide
    - Image embedding via file paths or base64
    """

    name = "pptx_generator"
    description = (
        "Generate polished PowerPoint presentations (.pptx) with slides containing "
        "text, images, tables, charts, and shapes. Includes 10 curated design "
        "palettes (e.g. Midnight Executive, Coral Energy, Ocean Gradient) with "
        "matched font pairings. Supports templates, gradient backgrounds, slide "
        "numbering, and speaker notes. Every slide should have a visual element — "
        "avoid text-only slides."
    )
    category = ToolCategory.GEN
    version = "2.0.0"
    timeout_sec = 120
    aliases = ["pptx_gen", "powerpoint_gen", "create_pptx", "presentation_gen"]
    search_hint = (
        "PowerPoint pptx generate create presentation slides deck layout "
        "theme palette design pitch report"
    )
    is_concurrency_safe = False
    is_read_only = False
    interrupt_behavior = "cancel"
    max_result_size_chars = 50_000
    path_params = ("template_path",)
    output_path_params = ("output_path",)

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "output_path": {
                    "type": "string",
                    "description": (
                        "Use a bare filename for the generated .pptx file "
                        "(for example, 'presentation.pptx'); it will be placed "
                        "in the session workspace and shown in the Files tab. "
                        "Use only when the user asked to save or export."
                    ),
                },
                "template_path": {
                    "type": "string",
                    "description": "Optional path to a .pptx template to use as base.",
                },
                "title": {
                    "type": "string",
                    "description": "Presentation title (metadata).",
                },
                "author": {
                    "type": "string",
                    "description": "Presentation author (metadata).",
                },
                "subject": {
                    "type": "string",
                    "description": "Presentation subject (metadata).",
                },
                "theme_palette": {
                    "type": "string",
                    "enum": list(THEME_PALETTES.keys()),
                    "description": (
                        "Curated color palette. Auto-applies primary/secondary/accent "
                        "colors and matched font pairing to all slides."
                    ),
                },
                "slide_width_inches": {
                    "type": "number",
                    "description": "Slide width in inches (default 13.333 for widescreen 16:9).",
                },
                "slide_height_inches": {
                    "type": "number",
                    "description": "Slide height in inches (default 7.5 for widescreen 16:9).",
                },
                "include_slide_numbers": {
                    "type": "boolean",
                    "description": "Add slide numbers to the bottom-right of each slide.",
                },
                "slides": {
                    "type": "array",
                    "description": "Array of slide definitions.",
                    "items": {
                        "type": "object",
                        "properties": {
                            "layout": {
                                "type": "string",
                                "enum": [
                                    "title",
                                    "title_content",
                                    "section_header",
                                    "two_content",
                                    "blank",
                                    "content_only",
                                    "title_only",
                                ],
                                "description": "Slide layout type.",
                            },
                            "title": {
                                "type": "string",
                                "description": "Slide title text.",
                            },
                            "subtitle": {
                                "type": "string",
                                "description": "Slide subtitle text (for title/section layouts).",
                            },
                            "content": {
                                "type": "array",
                                "description": "Content blocks to place on the slide.",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "type": {
                                            "type": "string",
                                            "enum": [
                                                "text",
                                                "bullet_list",
                                                "image",
                                                "table",
                                                "shape",
                                                "chart",
                                            ],
                                            "description": "Content block type.",
                                        },
                                        "text": {
                                            "type": "string",
                                            "description": "Text content.",
                                        },
                                        "items": {
                                            "type": "array",
                                            "items": {"type": "string"},
                                            "description": "Bullet list items.",
                                        },
                                        "level": {
                                            "type": "integer",
                                            "minimum": 0,
                                            "maximum": 4,
                                            "description": "Indentation level for bullets.",
                                        },
                                        "bold": {"type": "boolean"},
                                        "italic": {"type": "boolean"},
                                        "font_size": {
                                            "type": "integer",
                                            "description": "Font size in points.",
                                        },
                                        "font_name": {
                                            "type": "string",
                                            "description": (
                                                "Font family as installed on the machine that opens "
                                                "the deck. For mixed Chinese + English, use a "
                                                "pan-Unicode face or omit; names differ by Windows, "
                                                "macOS, and Linux."
                                            ),
                                        },
                                        "color": {
                                            "type": "string",
                                            "description": "Hex color (e.g. '#FF0000').",
                                        },
                                        "image_path": {
                                            "type": "string",
                                            "description": "Path to image file.",
                                        },
                                        "image_base64": {
                                            "type": "string",
                                            "description": "Base64-encoded image data.",
                                        },
                                        "image_mime": {
                                            "type": "string",
                                            "description": "MIME type for base64 image.",
                                        },
                                        "left": {"type": "number", "description": "Left position in inches."},
                                        "top": {"type": "number", "description": "Top position in inches."},
                                        "width": {"type": "number", "description": "Width in inches."},
                                        "height": {"type": "number", "description": "Height in inches."},
                                        "rows": {
                                            "type": "array",
                                            "description": "Table rows.",
                                            "items": {
                                                "type": "array",
                                                "items": {"type": "string"},
                                            },
                                        },
                                        "header_row": {
                                            "type": "boolean",
                                            "description": "First row is header.",
                                        },
                                        "chart_type": {
                                            "type": "string",
                                            "enum": ["bar", "line", "pie", "column"],
                                            "description": "Chart type.",
                                        },
                                        "chart_data": {
                                            "type": "object",
                                            "description": "Chart data with categories and series.",
                                        },
                                        "shape_type": {
                                            "type": "string",
                                            "enum": [
                                                "rectangle",
                                                "rounded_rectangle",
                                                "oval",
                                                "arrow_right",
                                                "arrow_left",
                                                "chevron",
                                            ],
                                        },
                                        "fill_color": {
                                            "type": "string",
                                            "description": "Shape fill color (hex).",
                                        },
                                    },
                                    "required": ["type"],
                                },
                            },
                            "notes": {
                                "type": "string",
                                "description": "Speaker notes for this slide.",
                            },
                            "background_color": {
                                "type": "string",
                                "description": "Slide background color (hex).",
                            },
                            "background_gradient": {
                                "type": "object",
                                "description": "Gradient background with two color stops.",
                                "properties": {
                                    "color1": {"type": "string", "description": "Start color (hex)."},
                                    "color2": {"type": "string", "description": "End color (hex)."},
                                },
                            },
                        },
                        "required": ["layout"],
                    },
                },
            },
            "required": ["output_path", "slides"],
            "additionalProperties": False,
        }

    def get_activity_description(self, params: dict[str, Any] | None = None) -> str | None:
        return "Generating PowerPoint presentation"

    def execute_sync(self, params: dict[str, Any], context: ToolContext) -> dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData
            from pptx.dml.color import RGBColor
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.util import Inches, Pt
        except ImportError as e:
            raise RuntimeError(
                "python-pptx is not installed. Install with: pip install python-pptx"
            ) from e

        output_path = Path(params["output_path"])
        template_path = params.get("template_path")
        slides_data = params.get("slides", [])
        palette_name = params.get("theme_palette")
        palette = THEME_PALETTES.get(palette_name or "", None)
        include_slide_numbers = params.get("include_slide_numbers", False)

        output_path.parent.mkdir(parents=True, exist_ok=True)

        if template_path:
            tmpl = Path(template_path)
            if not tmpl.exists():
                raise FileNotFoundError(f"Template not found: {tmpl}")
            prs = Presentation(str(tmpl))
        else:
            prs = Presentation()

        if params.get("slide_width_inches"):
            prs.slide_width = Inches(params["slide_width_inches"])
        if params.get("slide_height_inches"):
            prs.slide_height = Inches(params["slide_height_inches"])

        if params.get("title"):
            prs.core_properties.title = params["title"]
        if params.get("author"):
            prs.core_properties.author = params["author"]
        if params.get("subject"):
            prs.core_properties.subject = params["subject"]

        layout_map = self._build_layout_map(prs)

        stats = {"slides": 0, "text_blocks": 0, "images": 0, "tables": 0, "charts": 0, "shapes": 0}

        for slide_idx, slide_def in enumerate(slides_data):
            layout_name = slide_def.get("layout", "blank")
            slide_layout = layout_map.get(layout_name, layout_map.get("blank", prs.slide_layouts[6]))
            slide = prs.slides.add_slide(slide_layout)
            stats["slides"] += 1

            if slide_def.get("title") and slide.shapes.title:
                slide.shapes.title.text = slide_def["title"]
                if palette:
                    self._apply_palette_to_title(slide.shapes.title, palette, Pt, RGBColor)

            if slide_def.get("subtitle"):
                self._set_subtitle(slide, slide_def["subtitle"], palette, Pt, RGBColor)

            if slide_def.get("background_gradient"):
                self._set_gradient_background(
                    slide, slide_def["background_gradient"], RGBColor
                )
            elif slide_def.get("background_color"):
                self._set_slide_background(slide, slide_def["background_color"], RGBColor)

            for content_block in slide_def.get("content", []):
                block_type = content_block.get("type")
                self._apply_palette_defaults(content_block, palette)

                if block_type == "text":
                    self._add_text_block(slide, content_block, Inches, Pt, RGBColor)
                    stats["text_blocks"] += 1
                elif block_type == "bullet_list":
                    self._add_bullet_list(slide, content_block, Inches, Pt, RGBColor)
                    stats["text_blocks"] += 1
                elif block_type == "image":
                    self._add_image(slide, content_block, Inches)
                    stats["images"] += 1
                elif block_type == "table":
                    self._add_table(slide, content_block, Inches, Pt, RGBColor)
                    stats["tables"] += 1
                elif block_type == "chart":
                    self._add_chart(slide, content_block, Inches, CategoryChartData, XL_CHART_TYPE)
                    stats["charts"] += 1
                elif block_type == "shape":
                    self._add_shape(slide, content_block, Inches, Pt, RGBColor, MSO_SHAPE)
                    stats["shapes"] += 1

            if slide_def.get("notes"):
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_def["notes"]

            if include_slide_numbers:
                self._add_slide_number(slide, slide_idx + 1, prs, Inches, Pt, RGBColor)

        prs.save(str(output_path))
        file_size = output_path.stat().st_size

        logger.info(
            "pptx_generated",
            output_path=str(output_path),
            file_size=file_size,
            theme_palette=palette_name or "none",
            **stats,
        )

        return {
            "success": True,
            "output_path": str(output_path),
            "file_size_bytes": file_size,
            "content_stats": stats,
            "metadata": {
                "title": params.get("title", ""),
                "author": params.get("author", ""),
                "subject": params.get("subject", ""),
                "theme_palette": palette_name or "",
            },
        }

    # ------------------------------------------------------------------
    # Palette helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _apply_palette_defaults(block: dict[str, Any], palette: dict[str, str] | None) -> None:
        """Fill in font_name from palette when the block doesn't specify one."""
        if not palette:
            return
        if not block.get("font_name"):
            block.setdefault("_palette_font", palette.get("body_font"))

    @staticmethod
    def _apply_palette_to_title(
        title_shape: Any,
        palette: dict[str, str],
        Pt: Any,
        RGBColor: Any,
    ) -> None:
        """Style the title placeholder with palette colors and fonts."""
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = palette.get("header_font", "Calibri")
                run.font.size = Pt(TYPOGRAPHY_SIZES["slide_title"])
                try:
                    run.font.color.rgb = RGBColor.from_string(palette["primary"])
                except Exception:
                    pass

    def _set_subtitle(
        self,
        slide: Any,
        subtitle: str,
        palette: dict[str, str] | None,
        Pt: Any,
        RGBColor: Any,
    ) -> None:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = subtitle
                if palette:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = palette.get("body_font", "Calibri")
                            run.font.size = Pt(TYPOGRAPHY_SIZES["section_header"])
                return

    # ------------------------------------------------------------------
    # Slide backgrounds
    # ------------------------------------------------------------------

    @staticmethod
    def _set_slide_background(slide: Any, hex_color: str, RGBColor: Any) -> None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))

    @staticmethod
    def _set_gradient_background(
        slide: Any,
        gradient: dict[str, str],
        RGBColor: Any,
    ) -> None:
        """Apply a two-stop linear gradient background."""
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        color1 = gradient.get("color1", "FFFFFF")
        color2 = gradient.get("color2", "000000")
        stops = fill.gradient_stops
        stops[0].color.rgb = RGBColor.from_string(color1.lstrip("#"))
        stops[0].position = 0.0
        stops[1].color.rgb = RGBColor.from_string(color2.lstrip("#"))
        stops[1].position = 1.0

    # ------------------------------------------------------------------
    # Slide numbering
    # ------------------------------------------------------------------

    @staticmethod
    def _add_slide_number(
        slide: Any,
        num: int,
        prs: Any,
        Inches: Any,
        Pt: Any,
        RGBColor: Any,
    ) -> None:
        """Add a small slide-number text box to the bottom-right corner."""
        slide_w = prs.slide_width
        left = slide_w - Inches(0.8)
        top = prs.slide_height - Inches(0.45)
        txBox = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(num)
        p.alignment = 2  # right
        for run in p.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # ------------------------------------------------------------------
    # Layout map
    # ------------------------------------------------------------------

    @staticmethod
    def _build_layout_map(prs: Any) -> dict[str, Any]:
        """Map layout names to slide layout objects using index-based fallback."""
        layouts = prs.slide_layouts
        layout_map: dict[str, Any] = {}
        name_mapping = {
            "title": 0,
            "title_content": 1,
            "section_header": 2,
            "two_content": 3,
            "blank": 6,
            "content_only": 5,
            "title_only": 5,
        }
        for name, idx in name_mapping.items():
            if idx < len(layouts):
                layout_map[name] = layouts[idx]
        if "blank" not in layout_map and layouts:
            layout_map["blank"] = layouts[-1]
        return layout_map

    # ------------------------------------------------------------------
    # Content blocks
    # ------------------------------------------------------------------

    def _add_text_block(
        self, slide: Any, block: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any
    ) -> None:
        left = Inches(block.get("left", 0.5))
        top = Inches(block.get("top", 1.8))
        width = Inches(block.get("width", 8.0))
        height = Inches(block.get("height", 1.5))

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = block.get("text", "")
        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs:
            run.text = block.get("text", "")

        font_name = block.get("font_name") or block.get("_palette_font")
        if block.get("bold"):
            run.font.bold = True
        if block.get("italic"):
            run.font.italic = True
        if block.get("font_size"):
            run.font.size = Pt(block["font_size"])
        if font_name:
            run.font.name = font_name
        if block.get("color"):
            run.font.color.rgb = RGBColor.from_string(block["color"].lstrip("#"))

    def _add_bullet_list(
        self, slide: Any, block: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any
    ) -> None:
        left = Inches(block.get("left", 0.5))
        top = Inches(block.get("top", 1.8))
        width = Inches(block.get("width", 8.0))
        height = Inches(block.get("height", 4.0))

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        items = block.get("items", [])
        base_level = block.get("level", 0)
        font_name = block.get("font_name") or block.get("_palette_font")

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = base_level

            if block.get("font_size"):
                for run in p.runs:
                    run.font.size = Pt(block["font_size"])
            if font_name:
                for run in p.runs:
                    run.font.name = font_name
            if block.get("color"):
                for run in p.runs:
                    run.font.color.rgb = RGBColor.from_string(block["color"].lstrip("#"))

    def _add_image(self, slide: Any, block: dict[str, Any], Inches: Any) -> None:
        left = Inches(block.get("left", 0.5))
        top = Inches(block.get("top", 1.8))
        width = Inches(block.get("width", 5.0)) if block.get("width") else None
        height = Inches(block.get("height", 4.0)) if block.get("height") else None

        image_path = block.get("image_path")
        image_base64 = block.get("image_base64")

        if image_base64:
            image_data = base64.b64decode(image_base64)
            image_stream = io.BytesIO(image_data)
            if width and height:
                slide.shapes.add_picture(image_stream, left, top, width, height)
            elif width:
                slide.shapes.add_picture(image_stream, left, top, width=width)
            else:
                slide.shapes.add_picture(image_stream, left, top)
        elif image_path:
            img_file = Path(image_path)
            if not img_file.exists():
                logger.warning("pptx_image_not_found", image_path=image_path)
                return
            if width and height:
                slide.shapes.add_picture(str(img_file), left, top, width, height)
            elif width:
                slide.shapes.add_picture(str(img_file), left, top, width=width)
            else:
                slide.shapes.add_picture(str(img_file), left, top)

    def _add_table(
        self, slide: Any, block: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any
    ) -> None:
        rows_data = block.get("rows", [])
        if not rows_data:
            return

        num_rows = len(rows_data)
        num_cols = max(len(row) for row in rows_data)

        left = Inches(block.get("left", 0.5))
        top = Inches(block.get("top", 1.8))
        width = Inches(block.get("width", 8.0))
        height = Inches(block.get("height", 0.4 * num_rows + 0.4))

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        for i, row_data in enumerate(rows_data):
            for j, cell_text in enumerate(row_data):
                if j < num_cols:
                    table.cell(i, j).text = str(cell_text)

        if block.get("header_row") and num_rows > 0:
            for j in range(num_cols):
                cell = table.cell(0, j)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold = True

    def _add_chart(
        self, slide: Any, block: dict[str, Any], Inches: Any, CategoryChartData: Any, XL_CHART_TYPE: Any
    ) -> None:
        chart_type_str = block.get("chart_type", "column")
        chart_data_raw = block.get("chart_data", {})

        type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        xl_type = type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()
        categories = chart_data_raw.get("categories", [])
        chart_data.categories = categories

        for series in chart_data_raw.get("series", []):
            name = series.get("name", "Series")
            values = series.get("values", [])
            chart_data.add_series(name, values)

        left = Inches(block.get("left", 0.5))
        top = Inches(block.get("top", 1.8))
        width = Inches(block.get("width", 8.0))
        height = Inches(block.get("height", 4.5))

        slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)

    def _add_shape(
        self, slide: Any, block: dict[str, Any], Inches: Any, Pt: Any, RGBColor: Any, MSO_SHAPE: Any
    ) -> None:
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
            "chevron": MSO_SHAPE.CHEVRON,
        }
        shape_type = shape_map.get(block.get("shape_type", "rectangle"), MSO_SHAPE.RECTANGLE)

        left = Inches(block.get("left", 0.5))
        top = Inches(block.get("top", 1.8))
        width = Inches(block.get("width", 3.0))
        height = Inches(block.get("height", 2.0))

        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        if block.get("fill_color"):
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(block["fill_color"].lstrip("#"))

        if block.get("text"):
            shape.text = block["text"]
            font_name = block.get("font_name") or block.get("_palette_font")
            if block.get("font_size"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(block["font_size"])
            if font_name:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = font_name
            if block.get("color"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor.from_string(block["color"].lstrip("#"))
