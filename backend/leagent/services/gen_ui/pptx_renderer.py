"""Render a validated GenUI SlideDeck tree to a PowerPoint presentation (.pptx)."""

from __future__ import annotations

import io
from typing import Any


def render_genui_to_pptx(
    normalized_tree: dict[str, Any],
    *,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
) -> bytes:
    """Convert a normalized GenUI tree (typically a SlideDeck) into .pptx bytes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    root = normalized_tree.get("root")
    if not isinstance(root, dict):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        txBox.text_frame.text = "(empty)"
        return _save_to_bytes(prs)

    root_kind = str(root.get("kind") or "")

    if root_kind == "SlideDeck":
        slides = [
            c for c in (root.get("children") or [])
            if isinstance(c, dict)
        ]
        for slide_node in slides:
            _render_slide(prs, slide_node, Inches, Pt)
    elif root_kind == "Slide":
        _render_slide(prs, root, Inches, Pt)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _render_node_to_slide(slide, root, Inches, Pt, top_offset=1.0)

    return _save_to_bytes(prs)


def _save_to_bytes(prs: Any) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_slide(prs: Any, node: dict[str, Any], Inches: Any, Pt: Any) -> None:
    """Render a single Slide node into a new PowerPoint slide."""
    kind = str(node.get("kind") or "")
    props: dict[str, Any] = node.get("props") or {}
    children = [c for c in (node.get("children") or []) if isinstance(c, dict)]

    title = props.get("title", "")
    subtitle = props.get("subtitle", "")
    eyebrow = props.get("eyebrow", "")

    if title or subtitle:
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    else:
        slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    slide = prs.slides.add_slide(slide_layout)

    top_offset = 0.5

    if eyebrow:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(10), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = eyebrow.upper()
        run.font.size = Pt(11)
        run.font.bold = True
        top_offset += 0.5

    if title:
        if slide.shapes.title:
            slide.shapes.title.text = title
        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(11), Inches(1.0))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.size = Pt(28)
            run.font.bold = True
        top_offset += 1.2

    if subtitle:
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = subtitle
                break
        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(11), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = subtitle
            run.font.size = Pt(14)
        top_offset += 0.7

    for child in children:
        top_offset = _render_node_to_slide(slide, child, Inches, Pt, top_offset=top_offset)

    notes_text = props.get("notes", "")
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = str(notes_text)


def _render_node_to_slide(
    slide: Any, node: dict[str, Any], Inches: Any, Pt: Any, *, top_offset: float
) -> float:
    """Render a node onto an existing slide. Returns updated top_offset."""
    kind = str(node.get("kind") or "")
    props: dict[str, Any] = node.get("props") or {}
    children = [c for c in (node.get("children") or []) if isinstance(c, dict)]

    if kind == "Text":
        text = str(props.get("value", ""))
        if text:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(11), Inches(0.5))
            txBox.text_frame.word_wrap = True
            txBox.text_frame.paragraphs[0].text = text
            top_offset += 0.5

    elif kind == "Heading":
        level = int(props.get("level") or 2)
        text = str(props.get("value", ""))
        size = {1: 24, 2: 20, 3: 16}.get(level, 14)
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(11), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        top_offset += 0.7

    elif kind == "Table":
        headers = props.get("headers") or []
        rows_data: list[list[str]] = []
        if headers:
            rows_data.append([str(h) for h in headers])
        for child in children:
            if str(child.get("kind")) == "TableRow":
                row_cells = []
                for cell in (child.get("children") or []):
                    if isinstance(cell, dict):
                        cell_props = cell.get("props") or {}
                        row_cells.append(str(cell_props.get("value", "")))
                rows_data.append(row_cells)

        if rows_data:
            num_rows = len(rows_data)
            num_cols = max(len(r) for r in rows_data)
            table_shape = slide.shapes.add_table(
                num_rows, num_cols,
                Inches(0.8), Inches(top_offset),
                Inches(11), Inches(0.35 * num_rows + 0.3),
            )
            table = table_shape.table
            for i, row_data in enumerate(rows_data):
                for j, cell_text in enumerate(row_data):
                    if j < num_cols:
                        table.cell(i, j).text = cell_text
            top_offset += 0.35 * num_rows + 0.5

    elif kind == "Stat":
        value = str(props.get("value", ""))
        label = str(props.get("label", ""))
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(4), Inches(1.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.size = Pt(32)
        run.font.bold = True
        if label:
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = label
            run2.font.size = Pt(12)
        top_offset += 1.2

    elif kind == "KeyValueList":
        items = props.get("items") or []
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(top_offset), Inches(11), Inches(0.3 * len(items) + 0.3)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if not isinstance(item, dict):
                continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"{item.get('label', '')}: "
            run.font.bold = True
            p.add_run().text = str(item.get("value", ""))
        top_offset += 0.3 * len(items) + 0.4

    elif kind in ("Stack", "Row", "Grid", "ScrollArea", "Card", "DesignSurface", "AspectBox", "FeatureGrid"):
        for child in children:
            top_offset = _render_node_to_slide(slide, child, Inches, Pt, top_offset=top_offset)

    elif kind == "Image":
        src = props.get("src", "")
        if src:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(8), Inches(0.4))
            txBox.text_frame.paragraphs[0].text = f"[Image: {src}]"
            top_offset += 0.5

    elif kind == "Chart":
        title = str(props.get("title", "") or "")
        chart_t = str(props.get("chart", "line"))
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(11), Inches(0.45))
        txBox.text_frame.paragraphs[0].text = f"[Chart: {chart_t}] {title}".strip()
        top_offset += 0.55
        raw_cats = props.get("categories")
        categories = [str(x) for x in raw_cats] if isinstance(raw_cats, list) else []
        raw_series = props.get("series")
        rows_txt: list[str] = []
        if isinstance(raw_series, list):
            for s in raw_series:
                if not isinstance(s, dict):
                    continue
                nm = str(s.get("name", ""))
                vals = s.get("values") if isinstance(s.get("values"), list) else []
                rows_txt.append(f"{nm}: {vals}")
        if categories:
            rows_txt.insert(0, f"categories: {categories}")
        if rows_txt:
            tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(11), Inches(min(2.5, 0.35 * len(rows_txt))))
            tb2.text_frame.paragraphs[0].text = "\n".join(rows_txt)[:4000]
            top_offset += min(2.5, 0.35 * len(rows_txt)) + 0.2

    else:
        for child in children:
            top_offset = _render_node_to_slide(slide, child, Inches, Pt, top_offset=top_offset)
        if not children:
            text = str(props.get("value") or props.get("label") or props.get("title") or "")
            if text:
                txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top_offset), Inches(11), Inches(0.5))
                txBox.text_frame.paragraphs[0].text = text
                top_offset += 0.5

    return top_offset
