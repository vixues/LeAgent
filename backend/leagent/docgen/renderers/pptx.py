"""Professional PPTX renderer (python-pptx).

Slides are drawn on blank layouts with explicit theme-driven geometry — no
dependence on template layout indices. Composition (type scale, layout
regions, multi-level bullet plans, autofit) comes from
:mod:`leagent.docgen.slides`; this module is a thin translation of that
plan onto DrawingML shapes. Every text run gets an explicit east-asian
typeface (``<a:ea>``) so Chinese never falls back to a Latin-only face.

Content features:

- kicker (eyebrow) + action title + subtitle header band on every layout
- per-slide/deck backgrounds: solid, two-stop gradient, or cover-cropped
  image with a legibility scrim — text colors adapt to background darkness
- image-and-text layouts (``image.position``: right/left/top/full/background)
- 2-4 headed ``columns`` layouts with optional card emphasis
- leveled bullet typography with hanging indents, accent markers, numbered
  and task lists; body text auto-shrinks to fit its region
- bottom "takeaway" bar for the slide's so-what message
"""

from __future__ import annotations

import contextlib
import dataclasses
import io
from typing import TYPE_CHECKING, Any

import structlog

from leagent.docgen.charts import render_chart_png
from leagent.docgen.images import resolve_image
from leagent.docgen.markdown import parse_inline
from leagent.docgen.model import (
    DeckSpec,
    SlideBackground,
    SlideImage,
    SlideSpec,
    TableBlock,
)
from leagent.docgen.slides import (
    EMU_PER_PT,
    BulletPara,
    DeckTypography,
    Region,
    SlideGeometry,
    TextStyle,
    fit_body_size,
    flatten_body,
    is_dark_color,
    relative_luminance,
)
from leagent.docgen.tables import process_table, resolve_table_style
from leagent.docgen.themes import get_theme

if TYPE_CHECKING:
    from pathlib import Path

logger = structlog.get_logger(__name__)

_CHECK_COLOR = "#16A34A"


def render_pptx(deck: DeckSpec, output_path: Path) -> dict[str, Any]:
    """Render a deck spec to a .pptx file."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches, Pt

    theme = get_theme(deck.theme, kind="deck")
    typo = DeckTypography.from_theme(theme)
    warnings: list[str] = []
    stats = {"slides": 0, "images": 0, "charts": 0, "tables": 0}

    prs = Presentation()
    if deck.aspect == "16:9":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    geom = SlideGeometry(slide_w=int(prs.slide_width), slide_h=int(prs.slide_height))

    def _rgb(hex_color: str) -> Any:
        return RGBColor.from_string(hex_color.lstrip("#").upper()[:6])

    bg_color = _rgb(theme.colors.background)
    accent = _rgb(theme.colors.accent)
    surface = _rgb(theme.colors.surface)
    primary = _rgb(theme.colors.primary)
    on_primary = _rgb("#FFFFFF") if not theme.deck.dark else _rgb(theme.colors.text)

    # Role palettes. ``base`` matches the theme; the alternates keep text
    # readable when a custom slide background flips the effective darkness.
    _base_palette: dict[str, Any] = {
        "text": _rgb(theme.colors.text),
        "text_light": _rgb(theme.colors.text_light),
        # Dark decks: primary ≈ background, headings must use light text.
        "heading": _rgb(theme.colors.text) if theme.deck.dark else primary,
        "accent": accent,
        "on_primary": on_primary,
    }
    _dark_palette: dict[str, Any] = {
        "text": _rgb("#F0F2F5"),
        "text_light": _rgb("#C3C9D1"),
        "heading": _rgb("#FFFFFF"),
        "accent": accent,
        "on_primary": _rgb("#FFFFFF"),
    }
    _light_palette: dict[str, Any] = {
        "text": _rgb("#1F2933"),
        "text_light": _rgb("#5B6470"),
        "heading": primary,
        "accent": accent,
        "on_primary": _rgb("#FFFFFF"),
    }

    def _palette(dark: bool) -> dict[str, Any]:
        if dark == theme.deck.dark:
            return _base_palette
        return _dark_palette if dark else _light_palette

    # ------------------------------------------------------------------
    # Text primitives
    # ------------------------------------------------------------------

    def _style_run(
        run: Any,
        style: TextStyle,
        palette: dict[str, Any],
        *,
        scale: float = 1.0,
        bold: bool | None = None,
        italic: bool = False,
        mono: bool = False,
        color: Any = None,
    ) -> None:
        font_role = "mono" if mono else style.font
        run.font.name = getattr(theme.fonts, font_role)
        run.font.size = Pt(max(8.0, style.size * scale))
        run.font.bold = style.bold if bold is None else bold
        run.font.italic = style.italic or italic
        run.font.color.rgb = color if color is not None else palette[style.color]
        rpr = run._r.get_or_add_rPr()
        if style.letter_spacing_pt:
            rpr.set("spc", str(int(style.letter_spacing_pt * 100)))
        ea = rpr.find(qn("a:ea"))
        if ea is None:
            ea = rpr.makeelement(qn("a:ea"), {})
            rpr.append(ea)
        ea.set("typeface", theme.fonts.east_asia)

    def _fill_rich(
        paragraph: Any,
        text: str,
        style: TextStyle,
        palette: dict[str, Any],
        *,
        scale: float = 1.0,
        color: Any = None,
    ) -> None:
        if style.uppercase:
            text = text.upper()
        for span in parse_inline(text):
            if span.math:
                # Native OMML equation via mc:AlternateContent — PowerPoint
                # 2010+ renders the real equation; other viewers fall back to
                # the Unicode transliteration carried in the same element.
                from leagent.docgen.mathtext import latex_to_unicode
                from leagent.docgen.omml import omml_pptx_alternate

                unicode_text = latex_to_unicode(span.text)
                alt = omml_pptx_alternate(span.text, unicode_text, display=False)
                if alt is not None:
                    paragraph._p.append(alt)
                    continue
                run = paragraph.add_run()
                run.text = unicode_text
                _style_run(
                    run, style, palette, scale=scale, italic=True, color=color
                )
                continue
            run = paragraph.add_run()
            run.text = span.text
            _style_run(
                run,
                style,
                palette,
                scale=scale,
                bold=True if span.bold else None,
                italic=span.italic,
                mono=span.code,
                color=color,
            )
            if span.sup or span.sub:
                rpr = run._r.get_or_add_rPr()
                rpr.set("baseline", "30000" if span.sup else "-25000")
            if span.link:
                run.hyperlink.address = span.link

    def _para_spacing(paragraph: Any, style: TextStyle, *, scale: float = 1.0) -> None:
        paragraph.line_spacing = style.line_spacing
        if style.space_before_pt:
            paragraph.space_before = Pt(style.space_before_pt * scale)
        paragraph.space_after = Pt(style.space_after_pt * scale)

    def _set_indent(paragraph: Any, indent_emu: int, hang_emu: int) -> None:
        p_pr = paragraph._p.get_or_add_pPr()
        p_pr.set("marL", str(max(0, indent_emu + hang_emu)))
        p_pr.set("indent", str(-hang_emu if hang_emu else 0))

    def _set_bullet(
        paragraph: Any,
        *,
        char: str | None = None,
        auto_number: bool = False,
        color: Any = None,
    ) -> None:
        p_pr = paragraph._p.get_or_add_pPr()
        for tag in ("a:buClr", "a:buChar", "a:buNone", "a:buAutoNum", "a:buFont"):
            for el in p_pr.findall(qn(tag)):
                p_pr.remove(el)
        if char is None and not auto_number:
            p_pr.append(p_pr.makeelement(qn("a:buNone"), {}))
            return
        if color is not None:
            bu_clr = p_pr.makeelement(qn("a:buClr"), {})
            srgb = p_pr.makeelement(qn("a:srgbClr"), {"val": str(color)})
            bu_clr.append(srgb)
            p_pr.append(bu_clr)
        p_pr.append(p_pr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        if auto_number:
            p_pr.append(p_pr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"}))
        else:
            p_pr.append(p_pr.makeelement(qn("a:buChar"), {"char": char or "\u2022"}))

    def _textbox(slide: Any, region: Region) -> Any:
        box = slide.shapes.add_textbox(
            Emu(region.left), Emu(region.top), Emu(region.width), Emu(region.height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        return box

    def _one_line(
        slide: Any,
        region: Region,
        text: str,
        style: TextStyle,
        palette: dict[str, Any],
        *,
        align: Any = None,
        scale: float = 1.0,
        color: Any = None,
    ) -> Any:
        box = _textbox(slide, region)
        p = box.text_frame.paragraphs[0]
        if align is not None:
            p.alignment = align
        _para_spacing(p, style, scale=scale)
        _fill_rich(p, text, style, palette, scale=scale, color=color)
        return box

    # ------------------------------------------------------------------
    # Body engine: leveled paragraph plans -> a text frame
    # ------------------------------------------------------------------

    def _render_paras(
        tf: Any,
        paras: list[BulletPara],
        palette: dict[str, Any],
        *,
        scale: float = 1.0,
    ) -> None:
        first = True

        def _para() -> Any:
            nonlocal first
            if first:
                first = False
                return tf.paragraphs[0]
            return tf.add_paragraph()

        accent_hex = theme.colors.accent.lstrip("#").upper()[:6]
        for idx, item in enumerate(paras):
            p = _para()
            if item.kind in ("bullet", "numbered"):
                lvl = typo.level(item.level)
                style = lvl.text
                # Suppress extra leading space on the very first bullet.
                if idx == 0:
                    style = dataclasses.replace(style, space_before_pt=0.0)
                p.level = min(item.level, 4)
                _para_spacing(p, style, scale=scale)
                _set_indent(
                    p,
                    int(lvl.indent_in * 914_400 * scale),
                    int(lvl.hang_in * 914_400 * scale),
                )
                if item.kind == "numbered":
                    _set_bullet(p, auto_number=True, color=accent_hex)
                elif item.checked is not None:
                    _set_bullet(p, char=None)
                    glyph = "\u2713 " if item.checked else "\u25a1 "
                    run = p.add_run()
                    run.text = glyph
                    _style_run(
                        run,
                        style,
                        palette,
                        scale=scale,
                        bold=item.checked is True,
                        color=_rgb(_CHECK_COLOR) if item.checked else palette["text_light"],
                    )
                else:
                    _set_bullet(p, char=lvl.glyph, color=accent_hex)
                _fill_rich(p, item.text, style, palette, scale=scale)
            elif item.kind == "heading":
                _para_spacing(p, typo.run_in_heading, scale=scale)
                _set_bullet(p, char=None)
                _fill_rich(p, item.text, typo.run_in_heading, palette, scale=scale)
            elif item.kind == "quote":
                _para_spacing(p, typo.body, scale=scale)
                _set_bullet(p, char=None)
                _fill_rich(
                    p,
                    f"\u201c{item.text}\u201d",
                    typo.body,
                    palette,
                    scale=scale,
                    color=palette["text_light"],
                )
            elif item.kind == "code":
                _para_spacing(p, typo.code, scale=scale)
                _set_bullet(p, char=None)
                run = p.add_run()
                run.text = item.text
                _style_run(run, typo.code, palette, scale=scale, mono=True)
            else:  # para
                _para_spacing(p, typo.body, scale=scale)
                _set_bullet(p, char=None)
                _fill_rich(p, item.text, typo.body, palette, scale=scale)

    def _render_body(
        slide: Any,
        markdown_text: str,
        region: Region,
        palette: dict[str, Any],
    ) -> None:
        paras = flatten_body(markdown_text)
        if not paras:
            return
        scale = fit_body_size(paras, typo, region)
        box = _textbox(slide, region)
        _render_paras(box.text_frame, paras, palette, scale=scale)

    # ------------------------------------------------------------------
    # Backgrounds
    # ------------------------------------------------------------------

    def _set_shape_alpha(shape: Any, opacity: float) -> None:
        """Set fill opacity (0-1) on a solid-filled shape."""
        sp_pr = shape._element.spPr
        for srgb in sp_pr.iter(qn("a:srgbClr")):
            alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100_000))})
            srgb.append(alpha)
            break

    def _cover_crop(data: bytes, target_w: int, target_h: int) -> bytes | None:
        """Center-crop image bytes to the slide aspect ratio (PNG out)."""
        from PIL import Image as PILImage

        try:
            with PILImage.open(io.BytesIO(data)) as im:
                im = im.convert("RGB")
                iw, ih = im.size
                target_ratio = target_w / target_h
                if iw / ih > target_ratio:
                    new_w = int(ih * target_ratio)
                    x0 = (iw - new_w) // 2
                    im = im.crop((x0, 0, x0 + new_w, ih))
                else:
                    new_h = int(iw / target_ratio)
                    y0 = (ih - new_h) // 2
                    im = im.crop((0, y0, iw, y0 + new_h))
                buf = io.BytesIO()
                im.save(buf, format="PNG")
                return buf.getvalue()
        except Exception:  # noqa: BLE001
            return None

    def _image_luminance(data: bytes) -> float:
        from PIL import Image as PILImage

        try:
            with PILImage.open(io.BytesIO(data)) as im:
                gray = im.convert("L").resize((1, 1))
                return gray.getpixel((0, 0)) / 255.0
        except Exception:  # noqa: BLE001
            return 0.5

    def _solid_bg(slide: Any, color: Any) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _apply_background(slide: Any, bg: SlideBackground | None) -> bool:
        """Apply background config; returns whether the slide reads as dark."""
        if bg is None:
            _solid_bg(slide, bg_color)
            return theme.deck.dark

        if bg.gradient and len(bg.gradient) >= 2:
            try:
                fill = slide.background.fill
                fill.gradient()
                stops = fill.gradient_stops
                stops[0].color.rgb = _rgb(bg.gradient[0])
                stops[1].color.rgb = _rgb(bg.gradient[1])
                with contextlib.suppress(ValueError, AttributeError):
                    fill.gradient_angle = bg.gradient_angle
            except Exception:  # noqa: BLE001 — viewer-specific gradient quirks
                _solid_bg(slide, _rgb(bg.gradient[0]))
            lum = sum(relative_luminance(c) for c in bg.gradient[:2]) / 2
            return lum < 0.35

        if bg.image_path or bg.image_url or bg.image_base64:
            resolved = resolve_image(
                path=bg.image_path, url=bg.image_url, base64_data=bg.image_base64
            )
            if resolved is not None:
                cropped = _cover_crop(resolved.data, geom.slide_w, geom.slide_h) or resolved.data
                slide.shapes.add_picture(
                    io.BytesIO(cropped), 0, 0,
                    width=Emu(geom.slide_w), height=Emu(geom.slide_h),
                )
                overlay = bg.overlay if bg.overlay > 0 else 0.0
                if overlay > 0:
                    from pptx.enum.shapes import MSO_SHAPE

                    scrim = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, 0, 0, Emu(geom.slide_w), Emu(geom.slide_h)
                    )
                    scrim.fill.solid()
                    scrim.fill.fore_color.rgb = _rgb(bg.overlay_color)
                    scrim.line.fill.background()
                    scrim.shadow.inherit = False
                    _set_shape_alpha(scrim, overlay)
                img_lum = _image_luminance(resolved.data)
                eff_lum = img_lum * (1 - overlay) + relative_luminance(bg.overlay_color) * overlay
                return eff_lum < 0.45
            warnings.append("Background image could not be resolved")
            _solid_bg(slide, bg_color)
            return theme.deck.dark

        if bg.color:
            _solid_bg(slide, _rgb(bg.color))
            return is_dark_color(bg.color)

        _solid_bg(slide, bg_color)
        return theme.deck.dark

    # ------------------------------------------------------------------
    # Slide chrome (header band, accent bar, footer, takeaway)
    # ------------------------------------------------------------------

    def _accent_bar(slide: Any, region: Region) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(region.left), Emu(region.top), Emu(region.width), Emu(region.height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False

    def _slide_chrome(
        slide: Any, index: int, total: int, palette: dict[str, Any], *, skip: bool
    ) -> None:
        if skip:
            return
        if deck.footer_text:
            _one_line(slide, geom.footer_text_region(), deck.footer_text, typo.footer, palette)
        if deck.show_slide_numbers:
            _one_line(
                slide,
                geom.slide_number_region(),
                f"{index}/{total}",
                typo.footer,
                palette,
                align=PP_ALIGN.RIGHT,
            )

    def _slide_header(slide: Any, sl: SlideSpec, palette: dict[str, Any]) -> int:
        """Kicker + title + accent bar + subtitle band; returns content top."""
        has_kicker = bool(sl.kicker)
        if not sl.title and not has_kicker:
            return geom.margin
        if has_kicker:
            _one_line(slide, geom.kicker_region(), sl.kicker or "", typo.kicker, palette)
        if sl.title:
            _one_line(
                slide,
                geom.title_region(has_kicker=has_kicker),
                sl.title,
                typo.slide_title,
                palette,
            )
        _accent_bar(slide, geom.accent_bar_region(has_kicker=has_kicker))
        if sl.subtitle:
            _one_line(
                slide,
                geom.subtitle_region(has_kicker=has_kicker),
                sl.subtitle,
                typo.subtitle,
                palette,
            )
        return geom.content_top(
            has_title=True, has_kicker=has_kicker, has_subtitle=bool(sl.subtitle)
        )

    def _takeaway_bar(slide: Any, text: str) -> None:
        """Bottom so-what band: surface card with an accent left edge."""
        from pptx.enum.shapes import MSO_SHAPE

        region = geom.takeaway_region()
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(region.left), Emu(region.top), Emu(region.width), Emu(region.height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = surface
        card.line.fill.background()
        card.shadow.inherit = False
        _accent_bar(
            slide,
            Region(region.left, region.top, int(4 * EMU_PER_PT), region.height),
        )
        text_region = Region(
            region.left + int(0.22 * 914_400),
            region.top,
            region.width - int(0.34 * 914_400),
            region.height,
        )
        box = _textbox(slide, text_region)
        tf = box.text_frame
        from pptx.enum.text import MSO_ANCHOR

        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        # Cards keep theme colors regardless of custom slide backgrounds.
        _fill_rich(p, text, typo.takeaway, _base_palette)

    # ------------------------------------------------------------------
    # Media + table shapes
    # ------------------------------------------------------------------

    def _add_image_shape(slide: Any, data: bytes, region: Region) -> None:
        from PIL import Image as PILImage

        try:
            with PILImage.open(io.BytesIO(data)) as im:
                iw, ih = im.size
        except Exception:  # noqa: BLE001
            iw, ih = (1280, 720)
        scale = min(region.width / iw, region.height / ih)
        w = int(iw * scale)
        h = int(ih * scale)
        slide.shapes.add_picture(
            io.BytesIO(data),
            Emu(region.left + (region.width - w) // 2),
            Emu(region.top + (region.height - h) // 2),
            width=Emu(w),
            height=Emu(h),
        )

    def _place_slide_image(
        slide: Any, sl_image: SlideImage, region: Region, idx: int
    ) -> bool:
        resolved = resolve_image(
            path=sl_image.path, base64_data=sl_image.base64_data, url=sl_image.url
        )
        if resolved is None:
            warnings.append(f"Slide {idx}: image could not be resolved")
            return False
        img_region = region
        if sl_image.caption:
            cap_h = int(0.38 * 914_400)
            img_region = Region(
                region.left, region.top, region.width, max(0, region.height - cap_h)
            )
            _one_line(
                slide,
                Region(region.left, img_region.bottom, region.width, cap_h),
                sl_image.caption,
                typo.caption,
                _base_palette,
                align=PP_ALIGN.CENTER,
            )
        _add_image_shape(slide, resolved.data, img_region)
        stats["images"] += 1
        return True

    def _add_table_shape(slide: Any, block: TableBlock, region: Region) -> None:
        pt = process_table(block, theme=theme)
        if not pt.header and not pt.body:
            return
        ts = resolve_table_style(theme, pt.style, dark=theme.deck.dark)

        all_rows = ([pt.header] if pt.header else []) + pt.body
        n_rows = len(all_rows)
        row_h = int(0.42 * 914_400)
        table_h = min(region.height, row_h * n_rows)
        shape = slide.shapes.add_table(
            n_rows, pt.col_count,
            Emu(region.left), Emu(region.top), Emu(region.width), Emu(table_h),
        )
        table = shape.table

        for c_idx, frac in enumerate(pt.width_fractions()):
            table.columns[c_idx].width = Emu(int(region.width * frac))

        body_start = 1 if pt.header else 0
        header_fill = _rgb(ts.header_fill) if ts.header_fill else None
        zebra_fill = _rgb(ts.zebra_fill) if ts.zebra_fill else None
        total_fill = _rgb(ts.total_fill) if ts.total_fill else None
        pos_color = _rgb(ts.positive)
        neg_color = _rgb(ts.negative)
        header_text = _rgb(ts.header_text) if ts.header_fill else _base_palette["heading"]

        cell_style = TextStyle(size=max(11.0, theme.deck.body_size - 4))
        for r_idx, row in enumerate(all_rows):
            is_header = bool(pt.header) and r_idx == 0
            body_idx = r_idx - body_start
            is_total = not is_header and body_idx == pt.total_row_index
            for c_idx, pcell in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.fill.solid()
                if is_header and header_fill is not None:
                    cell.fill.fore_color.rgb = header_fill
                elif is_total and total_fill is not None:
                    cell.fill.fore_color.rgb = total_fill
                elif (
                    pt.zebra
                    and zebra_fill is not None
                    and not is_header
                    and body_idx % 2 == 1
                ):
                    cell.fill.fore_color.rgb = zebra_fill
                else:
                    cell.fill.fore_color.rgb = bg_color
                p = cell.text_frame.paragraphs[0]
                if pcell.align in ("center", "right"):
                    p.alignment = (
                        PP_ALIGN.CENTER if pcell.align == "center" else PP_ALIGN.RIGHT
                    )
                if is_header:
                    color = header_text
                elif pcell.polarity:
                    color = pos_color if pcell.polarity == "positive" else neg_color
                else:
                    color = _base_palette["text"]
                _fill_rich(
                    p,
                    pcell.text,
                    cell_style,
                    _base_palette,
                    color=color,
                )
                if is_header or pcell.bold:
                    for run in p.runs:
                        run.font.bold = True
        stats["tables"] += 1

    def _render_column(slide: Any, col: Any, region: Region, idx: int) -> None:
        pad = int(0.14 * 914_400)
        inner = region
        if col.emphasis:
            from pptx.enum.shapes import MSO_SHAPE

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(region.left), Emu(region.top), Emu(region.width), Emu(region.height),
            )
            with contextlib.suppress(Exception):  # keep corners subtle
                card.adjustments[0] = 0.045
            card.fill.solid()
            card.fill.fore_color.rgb = surface
            card.line.color.rgb = accent
            card.line.width = Pt(0.75)
            card.shadow.inherit = False
            inner = region.inset(pad)
        y = inner.top
        if col.heading:
            head_h = int(0.42 * 914_400)
            _one_line(
                slide,
                Region(inner.left, y, inner.width, head_h),
                col.heading,
                TextStyle(
                    size=theme.deck.body_size + 1.0, bold=True,
                    color="heading", font="heading", space_after_pt=2.0,
                ),
                _base_palette,
            )
            y += head_h + int(0.06 * 914_400)
        if col.image is not None:
            img_h = int(min(inner.bottom - y, inner.height * 0.45))
            if img_h > 0 and _place_slide_image(
                slide, col.image, Region(inner.left, y, inner.width, img_h), idx
            ):
                y += img_h + int(0.08 * 914_400)
        if col.body and inner.bottom - y > 0:
            _render_body(
                slide,
                col.body,
                Region(inner.left, y, inner.width, inner.bottom - y),
                _base_palette,
            )

    # ------------------------------------------------------------------
    # Slide construction
    # ------------------------------------------------------------------

    slides = list(deck.slides)
    if not slides or slides[0].layout != "title":
        slides.insert(
            0,
            SlideSpec(
                layout="title",
                title=deck.title or "Presentation",
                subtitle=deck.subtitle,
            ),
        )
    total = len(slides)

    for idx, sl in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        layout = sl.layout
        bg_spec = sl.background or deck.background
        # An `image.position == "background"` promotes the slide image.
        if (
            bg_spec is None
            and sl.image is not None
            and sl.image.position == "background"
        ):
            bg_spec = SlideBackground(
                image_path=sl.image.path,
                image_url=sl.image.url,
                image_base64=sl.image.base64_data,
                overlay=0.45,
            )
        if layout == "section" and bg_spec is None:
            # Section dividers keep their signature primary/surface band.
            section_fill = theme.colors.surface if theme.deck.dark else theme.colors.primary
            _solid_bg(slide, _rgb(section_fill))
            dark = is_dark_color(section_fill)
        else:
            dark = _apply_background(slide, bg_spec)
        palette = _palette(dark)
        has_takeaway = bool(sl.takeaway) and layout not in (
            "title", "section", "closing", "quote"
        )

        if layout in ("title", "closing"):
            band_top = geom.slide_h // 3
            if sl.kicker:
                _one_line(
                    slide,
                    Region(geom.margin, band_top - int(0.42 * 914_400), geom.content_w, int(0.34 * 914_400)),
                    sl.kicker,
                    typo.kicker,
                    palette,
                    align=PP_ALIGN.CENTER,
                )
            _one_line(
                slide,
                Region(geom.margin, band_top, geom.content_w, int(1.6 * 914_400)),
                sl.title or (deck.title if layout == "title" else "Thank You"),
                typo.display,
                palette,
                align=PP_ALIGN.CENTER,
            )
            _accent_bar(
                slide,
                Region(
                    (geom.slide_w - int(2 * 914_400)) // 2,
                    band_top + int(1.5 * 914_400),
                    int(2 * 914_400),
                    int(4 * EMU_PER_PT),
                ),
            )
            sub_text = sl.subtitle or sl.body
            if layout == "title" and not sub_text:
                bits = [b for b in (deck.author, deck.date) if b]
                sub_text = "  ·  ".join(bits) if bits else None
            if sub_text:
                _one_line(
                    slide,
                    Region(
                        geom.margin,
                        band_top + int(1.75 * 914_400),
                        geom.content_w,
                        int(0.8 * 914_400),
                    ),
                    sub_text,
                    typo.subtitle,
                    palette,
                    align=PP_ALIGN.CENTER,
                )
            _slide_chrome(slide, idx, total, palette, skip=True)

        elif layout == "section":
            sect_palette = dict(palette)
            if sl.background is None and deck.background is None:
                sect_palette["heading"] = on_primary
                sect_palette["text"] = on_primary
                sect_palette["text_light"] = on_primary
            if sl.kicker:
                _one_line(
                    slide,
                    Region(
                        geom.margin,
                        geom.slide_h // 2 - int(1.24 * 914_400),
                        geom.content_w,
                        int(0.34 * 914_400),
                    ),
                    sl.kicker,
                    typo.kicker,
                    sect_palette,
                )
            _one_line(
                slide,
                Region(
                    geom.margin,
                    geom.slide_h // 2 - int(0.8 * 914_400),
                    geom.content_w,
                    int(1.4 * 914_400),
                ),
                sl.title or "",
                typo.section,
                sect_palette,
                color=sect_palette["heading"],
            )
            if sl.subtitle or sl.body:
                _one_line(
                    slide,
                    Region(
                        geom.margin,
                        geom.slide_h // 2 + int(0.6 * 914_400),
                        geom.content_w,
                        int(0.8 * 914_400),
                    ),
                    sl.subtitle or sl.body or "",
                    typo.subtitle,
                    sect_palette,
                    color=sect_palette["text"],
                )
            _slide_chrome(slide, idx, total, palette, skip=True)

        elif layout == "quote":
            region = Region(
                int(1.2 * 914_400),
                geom.slide_h // 3,
                geom.slide_w - int(2.4 * 914_400),
                int(2.2 * 914_400),
            )
            box = _textbox(slide, region)
            p = box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _para_spacing(p, typo.quote)
            _fill_rich(p, f"\u201c{sl.quote or sl.body or ''}\u201d", typo.quote, palette)
            if sl.attribution:
                ap = box.text_frame.add_paragraph()
                ap.alignment = PP_ALIGN.CENTER
                _para_spacing(ap, typo.attribution)
                _fill_rich(ap, f"— {sl.attribution}", typo.attribution, palette)
            _slide_chrome(slide, idx, total, palette, skip=False)

        elif layout == "two_column":
            top = _slide_header(slide, sl, palette)
            region = geom.content_region(top, has_takeaway=has_takeaway)
            cols = geom.columns(region, 2)
            for col_region, text in zip(cols, (sl.left, sl.right), strict=False):
                if text:
                    _render_body(slide, text, col_region, palette)
            if has_takeaway:
                _takeaway_bar(slide, sl.takeaway or "")
            _slide_chrome(slide, idx, total, palette, skip=False)

        elif layout == "columns":
            top = _slide_header(slide, sl, palette)
            region = geom.content_region(top, has_takeaway=has_takeaway)
            col_specs = sl.columns or []
            if col_specs:
                weights = (
                    [c.width or 1.0 for c in col_specs]
                    if any(c.width for c in col_specs)
                    else None
                )
                col_regions = geom.columns(region, len(col_specs), weights=weights)
                for col, col_region in zip(col_specs, col_regions, strict=False):
                    _render_column(slide, col, col_region, idx)
            elif sl.body:
                _render_body(slide, sl.body, region, palette)
            if has_takeaway:
                _takeaway_bar(slide, sl.takeaway or "")
            _slide_chrome(slide, idx, total, palette, skip=False)

        elif layout == "image":
            top = _slide_header(slide, sl, palette)
            region = geom.content_region(top, has_takeaway=has_takeaway)
            position = sl.image.position if sl.image else "full"
            placed = False
            if sl.image and position == "background":
                # Image already applied as slide background; body gets the
                # full content region on top of the scrim.
                if sl.body:
                    _render_body(slide, sl.body, region, palette)
                placed = True
            elif sl.image and position in ("left", "right", "top") and sl.body:
                text_region, media_region = geom.split(
                    region, sl.image.ratio, side=position
                )
                placed = _place_slide_image(slide, sl.image, media_region, idx)
                _render_body(slide, sl.body, text_region, palette)
            elif sl.image:
                placed = _place_slide_image(slide, sl.image, region, idx)
            if not placed and not sl.body and sl.image is None:
                warnings.append(f"Slide {idx}: image layout without image")
            if not placed and sl.body and position not in ("left", "right", "top"):
                _render_body(slide, sl.body, region, palette)
            if has_takeaway:
                _takeaway_bar(slide, sl.takeaway or "")
            _slide_chrome(slide, idx, total, palette, skip=False)

        elif layout == "table":
            top = _slide_header(slide, sl, palette)
            region = geom.content_region(top, has_takeaway=has_takeaway)
            if sl.table:
                _add_table_shape(slide, sl.table, region)
            elif sl.body:
                _render_body(slide, sl.body, region, palette)
            if has_takeaway:
                _takeaway_bar(slide, sl.takeaway or "")
            _slide_chrome(slide, idx, total, palette, skip=False)

        elif layout == "chart":
            top = _slide_header(slide, sl, palette)
            region = geom.content_region(top, has_takeaway=has_takeaway)
            chart_block = sl.chart
            if chart_block is not None:
                png = render_chart_png(
                    chart_block,
                    theme,
                    width_in=9.0,
                    height_in=4.6,
                    transparent=not theme.deck.dark,
                )
                if png is not None:
                    if sl.body:
                        text_region, media_region = geom.split(region, 0.58, side="left")
                        _render_body(slide, sl.body, text_region, palette)
                        _add_image_shape(slide, png, media_region)
                    else:
                        _add_image_shape(slide, png, region)
                    stats["charts"] += 1
                else:
                    warnings.append(f"Slide {idx}: chart could not be rendered")
            elif sl.body:
                _render_body(slide, sl.body, region, palette)
            if has_takeaway:
                _takeaway_bar(slide, sl.takeaway or "")
            _slide_chrome(slide, idx, total, palette, skip=False)

        else:  # content
            top = _slide_header(slide, sl, palette)
            region = geom.content_region(top, has_takeaway=has_takeaway)
            if sl.image is not None and sl.image.position != "background":
                side = sl.image.position if sl.image.position in ("left", "right", "top") else "right"
                text_region, media_region = geom.split(region, sl.image.ratio, side=side)
                _place_slide_image(slide, sl.image, media_region, idx)
                if sl.body:
                    _render_body(slide, sl.body, text_region, palette)
            elif sl.body and sl.table:
                # split(side="top") returns (lower, upper): body above, table below.
                lower_region, upper_region = geom.split(region, 0.5, side="top")
                _render_body(slide, sl.body, upper_region, palette)
                _add_table_shape(slide, sl.table, lower_region)
            elif sl.table:
                _add_table_shape(slide, sl.table, region)
            elif sl.body:
                _render_body(slide, sl.body, region, palette)
            if has_takeaway:
                _takeaway_bar(slide, sl.takeaway or "")
            _slide_chrome(slide, idx, total, palette, skip=False)

        if sl.notes:
            slide.notes_slide.notes_text_frame.text = sl.notes
        stats["slides"] += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    result = {
        "success": True,
        "output_path": str(output_path),
        "format": "pptx",
        "file_size_bytes": output_path.stat().st_size,
        "slide_count": stats["slides"],
        "content_stats": stats,
        "font_embedded": False,  # OOXML text relies on viewer fonts by design
        "east_asia_font": theme.fonts.east_asia,
        "theme": theme.name,
        "warnings": warnings,
    }
    logger.info("docgen_pptx_rendered", output_path=str(output_path), **stats)
    return result
